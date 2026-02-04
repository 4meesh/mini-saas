from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()

    # Layouts: 0 is Title, 1 is Title + Content
    
    # 1. Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Mini SaaS Task Manager"
    subtitle.text = "A Modern, Productive, and Beautiful Task Management Solution\nBuilt with React & Supabase"
    
    # 2. Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Overview"
    content = slide.placeholders[1]
    content.text = (
        "The Mini SaaS Task Manager is a full-stack application designed for simplicity and productivity. "
        "It features a real-time dashboard, interactive task management, and advanced data portability options."
    )
    
    # 3. Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Features"
    tf = slide.placeholders[1].text_frame
    tf.text = "Comprehensive Task Control"
    
    p = tf.add_paragraph()
    p.text = "• CRUD Operations: Create, Edit, Toggle, and Delete tasks in real-time."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Advanced Categorization: Priority labeling and situational categories."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Data Portability: Export to JSON/CSV and Import from JSON."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Themes: Fluid transitions between Light and Amoled Dark modes."
    p.level = 1

    # 4. Tech Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Technical Architecture"
    tf = slide.placeholders[1].text_frame
    
    p = tf.add_paragraph()
    p.text = "Frontend Framework: React.js (Vite)"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Backend & Auth: Supabase (PostgreSQL)"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Styling: Vanilla CSS with Modern Variables & Glassmorphism"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Real-time State: Supabase Auth & DB Listeners"
    p.font.bold = True

    # 5. Visual Excellence (Light Mode)
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # 5 is Title Only
    slide.shapes.title.text = "Modern Light Interface"
    
    img_path = r'C:/Users/Ameesh Mohammed/.gemini/antigravity/brain/dafd483c-8b2e-40fc-b0ba-db937d09bf06/saas_dashboard_light_mockup_1770196765969.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(5))
    
    # 6. Visual Excellence (Dark Mode)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Ultra Dark 'Amoled' Mode"
    
    img_path = r'C:/Users/Ameesh Mohammed/.gemini/antigravity/brain/dafd483c-8b2e-40fc-b0ba-db937d09bf06/saas_dashboard_dark_mockup_1770196796593.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(5))

    # 7. Statistics & Insights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dashboard Analytics"
    tf = slide.placeholders[1].text_frame
    tf.text = "The dashboard provides instant visual feedback:"
    p = tf.add_paragraph()
    p.text = "• Real-time Task Breakdown by Priority."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Completion Rate Progress Visualization."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Category-wise Active Task Tracking."
    p.level = 1

    # 8. Success & Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You!"
    subtitle.text = "A Premium Task Management Experience\nQuestions & Feedback"

    prs.save('Mini_SaaS_Project_Presentation.pptx')
    print("Presentation created successfully: Mini_SaaS_Project_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
